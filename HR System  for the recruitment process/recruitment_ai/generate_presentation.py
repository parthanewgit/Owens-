#!/usr/bin/env python3
"""
Recruitment AI System - PowerPoint Presentation Generator
Generates a professional, colorful presentation with all content
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import os

# Color Scheme
PRIMARY_BLUE = RGBColor(0, 109, 171)      # #006DAB
ACCENT_GREEN = RGBColor(0, 166, 81)       # #00A651
DARK_GRAY = RGBColor(44, 62, 80)          # #2C3E50
LIGHT_GRAY = RGBColor(236, 240, 241)      # #ECF0F1
WHITE = RGBColor(255, 255, 255)           # #FFFFFF

class PresentationGenerator:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)

    def add_title_slide(self, title, subtitle, presenter):
        """Add a title slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank layout
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = PRIMARY_BLUE

        # Add gradient effect with shapes
        shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(7.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = PRIMARY_BLUE
        shape.line.color.rgb = PRIMARY_BLUE

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(60)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = WHITE
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Subtitle
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(9), Inches(1))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        subtitle_frame.paragraphs[0].font.size = Pt(28)
        subtitle_frame.paragraphs[0].font.color.rgb = ACCENT_GREEN
        subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Presenter
        presenter_box = slide.shapes.add_textbox(Inches(0.5), Inches(6), Inches(9), Inches(1))
        presenter_frame = presenter_box.text_frame
        presenter_frame.text = presenter
        presenter_frame.paragraphs[0].font.size = Pt(16)
        presenter_frame.paragraphs[0].font.color.rgb = WHITE
        presenter_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    def add_content_slide(self, title, content_left, content_right=None):
        """Add a content slide with title and bullet points"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = WHITE

        # Title bar
        title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
        title_shape.fill.solid()
        title_shape.fill.fore_color.rgb = PRIMARY_BLUE
        title_shape.line.color.rgb = PRIMARY_BLUE

        # Title text
        title_frame = title_shape.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(40)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = WHITE
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        title_frame.margin_top = Inches(0.1)

        # Content
        if content_right is None:
            # Single column
            content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8.4), Inches(6))
            text_frame = content_box.text_frame
            text_frame.word_wrap = True
            
            for i, point in enumerate(content_left):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                p.text = point
                p.level = 0
                p.font.size = Pt(18)
                p.font.color.rgb = DARK_GRAY
                p.space_before = Pt(8)
                p.space_after = Pt(8)
        else:
            # Two columns
            left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.5), Inches(6))
            left_frame = left_box.text_frame
            left_frame.word_wrap = True
            
            for i, point in enumerate(content_left):
                if i == 0:
                    p = left_frame.paragraphs[0]
                else:
                    p = left_frame.add_paragraph()
                p.text = point
                p.font.size = Pt(16)
                p.font.color.rgb = DARK_GRAY
                p.space_before = Pt(6)

            right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.2), Inches(4.3), Inches(6))
            right_frame = right_box.text_frame
            right_frame.word_wrap = True
            
            for i, point in enumerate(content_right):
                if i == 0:
                    p = right_frame.paragraphs[0]
                else:
                    p = right_frame.add_paragraph()
                p.text = point
                p.font.size = Pt(16)
                p.font.color.rgb = DARK_GRAY
                p.space_before = Pt(6)

    def add_metric_slide(self, title, metrics):
        """Add a slide with key metrics/statistics"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = WHITE

        # Title bar
        title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
        title_shape.fill.solid()
        title_shape.fill.fore_color.rgb = PRIMARY_BLUE
        title_shape.line.color.rgb = PRIMARY_BLUE
        
        title_frame = title_shape.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(40)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = WHITE
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Add metric boxes
        box_width = 4.5
        box_height = 2
        start_x = 0.8
        start_y = 1.5

        for idx, (label, value) in enumerate(metrics):
            col = idx % 2
            row = idx // 2
            left = start_x + col * (box_width + 0.5)
            top = start_y + row * (box_height + 0.5)

            # Box background
            box = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(box_width), Inches(box_height))
            box.fill.solid()
            box.fill.fore_color.rgb = LIGHT_GRAY
            box.line.color.rgb = PRIMARY_BLUE
            box.line.width = Pt(2)

            # Value
            value_box = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.2), Inches(box_width - 0.4), Inches(0.8))
            value_frame = value_box.text_frame
            value_frame.text = value
            value_frame.paragraphs[0].font.size = Pt(32)
            value_frame.paragraphs[0].font.bold = True
            value_frame.paragraphs[0].font.color.rgb = ACCENT_GREEN
            value_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

            # Label
            label_box = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 1), Inches(box_width - 0.4), Inches(0.8))
            label_frame = label_box.text_frame
            label_frame.text = label
            label_frame.word_wrap = True
            label_frame.paragraphs[0].font.size = Pt(14)
            label_frame.paragraphs[0].font.color.rgb = DARK_GRAY
            label_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    def generate(self):
        """Generate the complete presentation"""
        
        # SLIDE 1: Title Slide
        self.add_title_slide(
            "RECRUITMENT AI SYSTEM",
            "Transforming Hiring Through Intelligent Automation & AI",
            "Parthasarathi Mishra | Prateek Anand | March 24, 2026"
        )

        # SLIDE 2: Problem & Solution
        self.add_content_slide(
            "Problem Statement & Solution",
            [
                "❌ TRADITIONAL HIRING CHALLENGES",
                "• Time-consuming (2-8 hours per role)",
                "• Subjective evaluation",
                "• Inconsistent processes",
                "• Manual skill gap analysis",
                "• High hiring costs",
            ],
            [
                "✅ AI-POWERED SOLUTION",
                "• Automated resume analysis (seconds)",
                "• Objective AI-based scoring",
                "• Standardized assessment",
                "• Real-time skill gap detection",
                "• 60-70% cost reduction",
            ]
        )

        # SLIDE 3: System Overview
        self.add_metric_slide(
            "System Overview - Key Capabilities",
            [
                ("Resume Analysis", "3-4 sec"),
                ("Accuracy Rate", "92%+"),
                ("Question Generation", "5-7 sec"),
                ("Concurrent Users", "100+"),
            ]
        )

        # SLIDE 4: Architecture Overview
        self.add_content_slide(
            "System Architecture - High Level",
            [
                "🎨 FRONTEND LAYER",
                "Streamlit Web Application (Port 8501)",
                "Interactive dashboard for users",
                "",
                "⚙️ API GATEWAY",
                "FastAPI Backend (Port 8000)",
                "CORS, Authentication, Routing",
            ],
            [
                "🤖 AI AGENTS LAYER",
                "Resume Agent - Skill matching",
                "Question Agent - Question generation",
                "Evaluation Agent - Answer scoring",
                "Learning Agent - Roadmap creation",
                "",
                "🧠 EXTERNAL AI",
                "OpenAI GPT-4o + Embeddings",
            ]
        )

        # SLIDE 5: Core Agents - Resume Agent
        self.add_content_slide(
            "Core Agent #1: Resume Matching Agent",
            [
                "📄 FUNCTIONALITY",
                "✓ Skill extraction from resumes",
                "✓ Semantic matching with JD",
                "✓ Candidate ranking",
                "✓ Skill gap identification",
                "",
                "📊 OUTPUT",
                "Ranked candidates with scores",
                "Match percentages (0-100%)",
                "Skill analysis per candidate",
            ],
            [
                "⚡ PERFORMANCE",
                "Processing Time: 3-4 sec",
                "Accuracy: 92%+",
                "",
                "🔄 WORKFLOW",
                "1. Parse resume",
                "2. Extract skills",
                "3. Create embeddings",
                "4. Calculate similarity",
                "5. Rank candidates",
                "6. Return results",
            ]
        )

        # SLIDE 6: Core Agents - Question Agent
        self.add_content_slide(
            "Core Agent #2: Question Generation Agent",
            [
                "❓ FUNCTIONALITY",
                "✓ Generate technical questions",
                "✓ Create behavioral questions",
                "✓ Develop scenario questions",
                "✓ Customize by role",
                "✓ Provide expected answers",
                "",
                "📋 QUESTION TYPES",
                "Technical: 5 questions",
                "Behavioral: 3 questions",
                "Scenario: 2 questions",
            ],
            [
                "⚡ PERFORMANCE",
                "Processing Time: 5-7 sec",
                "Total Questions per role: 10",
                "",
                "📝 OUTPUT FORMATS",
                "Multiple Choice (MCQ)",
                "Free-text questions",
                "Expected answers",
                "Rubrics included",
                "",
                "🎯 CUSTOMIZATION",
                "Role-specific questions",
                "Difficulty levels",
            ]
        )

        # SLIDE 7: Core Agents - Evaluation Agent
        self.add_content_slide(
            "Core Agent #3: Answer Evaluation Agent",
            [
                "✅ SCORING METHODOLOGY",
                "MCQ: Instant scoring (0 or 10)",
                "Free-text: Semantic analysis",
                "Overall: Average of scores",
                "Scale: 0-100 points",
                "",
                "📊 OUTPUT METRICS",
                "Total score (0-10)",
                "Per-question breakdown",
                "Detailed feedback",
                "Strengths & improvements",
            ],
            [
                "⚡ PERFORMANCE",
                "MCQ: 2-3 sec",
                "Free-text: 3-5 sec",
                "Overall accuracy: 95%+",
                "",
                "💡 FEATURES",
                "✓ Natural language understanding",
                "✓ Context-aware evaluation",
                "✓ Constructive feedback",
                "✓ Comparative analytics",
                "✓ Performance tracking",
            ]
        )

        # SLIDE 8: Core Agents - Learning Agent
        self.add_content_slide(
            "Core Agent #4: Learning Roadmap Agent",
            [
                "📚 FUNCTIONALITY",
                "✓ Skill gap analysis",
                "✓ Personalized learning paths",
                "✓ Resource recommendations",
                "✓ Timeline estimation",
                "✓ Progress tracking setup",
                "",
                "🎯 OUTPUT COMPONENTS",
                "Missing skills list",
                "Learning objectives",
                "Recommended resources",
                "Estimated duration",
            ],
            [
                "⚡ PERFORMANCE",
                "Processing Time: 4-5 sec",
                "Customization level: High",
                "",
                "✨ PERSONALIZATION",
                "Individual skill assessment",
                "Learning style compatibility",
                "Role-specific paths",
                "Progressive milestones",
                "Success metrics",
                "",
                "📈 TRACKING",
                "Progress monitoring",
            ]
        )

        # SLIDE 9: Complete Workflow
        self.add_content_slide(
            "Complete User Workflow",
            [
                "STEP 1: LOGIN",
                "→ Enter credentials (admin/admin123)",
                "→ Receive JWT token",
                "",
                "STEP 2: RESUME ANALYSIS",
                "→ Upload job description",
                "→ Upload resumes (PDF/DOCX)",
                "→ View ranked candidates (92% accuracy)",
                "",
                "STEP 3: INTERVIEW PREP",
                "→ System generates 10 tailored questions",
                "→ 5 technical + 3 behavioral + 2 scenario",
            ],
            [
                "STEP 4: ANSWER EVALUATION",
                "→ Candidate submits answers",
                "→ AI evaluates instantly",
                "→ Receive score (0-10) + feedback",
                "",
                "STEP 5: LEARNING PATH",
                "→ Review skill gaps",
                "→ Get personalized learning roadmap",
                "→ Access timeline & resources",
                "",
                "TIMELINE",
                "Complete assessment: ~30 min",
                "Results available: Immediately",
            ]
        )

        # SLIDE 10: Technology Stack
        self.add_content_slide(
            "Technology Stack",
            [
                "🎨 FRONTEND",
                "Streamlit 1.28.1 - Web UI",
                "Requests 2.31.0 - API calls",
                "",
                "⚙️ BACKEND",
                "FastAPI 0.104.1 - API Framework",
                "Uvicorn 0.24.0 - ASGI Server",
                "Pydantic 2.4.2 - Validation",
                "",
                "🧠 AI & PROCESSING",
                "OpenAI GPT-4o - Language Model",
                "OpenAI Embeddings - Vector rep",
                "scikit-learn 1.3.2 - ML algorithms",
            ],
            [
                "💾 DATA & STORAGE",
                "ChromaDB 0.4.14 - Vector store",
                "SQLite - Relational DB",
                "SQLAlchemy 2.0.23 - ORM",
                "",
                "🔐 SECURITY",
                "PyJWT 2.11.0 - Token management",
                "python-multipart - Form handling",
                "",
                "📄 DOCUMENT PROCESSING",
                "PyPDF2 3.0.1 - PDF parsing",
                "python-docx 0.8.11 - DOCX parsing",
                "",
                "✅ TESTING",
                "pytest 7.4.3 - Testing framework",
            ]
        )

        # SLIDE 11: Architecture Diagram
        self.add_content_slide(
            "System Architecture Diagram",
            [
                "┌─────────────────────────┐",
                "│   STREAMLIT FRONTEND    │",
                "│   (Port 8501)           │",
                "└───────────┬─────────────┘",
                "            │ HTTP",
                "┌───────────▼─────────────┐",
                "│  FASTAPI BACKEND        │",
                "│  (Port 8000)            │",
                "└───────────┬─────────────┘",
                "    ┌───────┴───────┐",
                "    ▼               ▼",
                "┌────────┐     ┌──────────┐",
                "│ Agents │     │OpenAI    │",
                "└────────┘     │GPT-4o    │",
                "               └──────────┘",
                "",
                "┌────────────────────────┐",
                "│ SQLite │  ChromaDB    │",
                "│ (DB)   │  (Vectors)   │",
                "└────────────────────────┘",
            ]
        )

        # SLIDE 12: Security & Authentication
        self.add_content_slide(
            "Security & Authentication",
            [
                "🔐 JWT TOKEN FLOW",
                "1. User login → Credentials",
                "2. Verify against user DB",
                "3. Generate JWT token",
                "4. 24-hour expiration",
                "5. Include in API requests",
                "",
                "🛡️ SECURITY FEATURES",
                "✓ Password hashing",
                "✓ Token-based auth",
                "✓ CORS protection",
                "✓ Input validation",
                "✓ SQL parameterization",
            ],
            [
                "🔑 TOKEN STRUCTURE",
                "Header: {alg: HS256, typ: JWT}",
                "Payload: {sub, exp, iat}",
                "Signature: HMAC-SHA256",
                "",
                "📋 API PROTECTION",
                "All endpoints require JWT",
                "Bearer token in header",
                "40-minute request timeout",
                "",
                "✅ COMPLIANCE",
                "Role-based access control",
                "Audit logging enabled",
                "No sensitive data exposure",
            ]
        )

        # SLIDE 13: Performance Metrics
        self.add_metric_slide(
            "System Performance Metrics",
            [
                ("Resume Analysis", "3-4 sec"),
                ("Question Generation", "5-7 sec"),
                ("Answer Evaluation", "2-3 sec"),
                ("Accuracy Rate", "92-95%"),
            ]
        )

        # SLIDE 14: Key Features & Benefits
        self.add_content_slide(
            "Key Features & Business Benefits",
            [
                "⚡ SPEED & EFFICIENCY",
                "Resume analysis in seconds",
                "60-70% reduction in hiring time",
                "Instant question generation",
                "Real-time answer evaluation",
                "",
                "🎯 OBJECTIVITY & FAIRNESS",
                "AI-based objective scoring",
                "Eliminates human bias",
                "Consistent evaluation",
                "Standardized process",
                "",
                "💡 INTELLIGENCE",
                "Semantic skill matching",
                "Context-aware evaluation",
                "Personalized recommendations",
            ],
            [
                "💰 COST REDUCTION",
                "90% reduction in screening time",
                "40% cost per hire",
                "ROI within 6 months",
                "Scalable solution",
                "",
                "📊 ADVANCED ANALYTICS",
                "Real-time hiring metrics",
                "Data-driven decisions",
                "Candidate comparisons",
                "Performance tracking",
                "",
                "🔄 CONTINUOUS IMPROVEMENT",
                "Learning from feedback",
                "Model improvement",
                "Enhanced accuracy",
            ]
        )

        # SLIDE 15: Use Cases
        self.add_content_slide(
            "Real-World Use Cases",
            [
                "🏢 LARGE ENTERPRISE",
                "Scenario: Hiring 50 developers",
                "Timeline: 8 weeks → 2-3 weeks",
                "Cost reduction: 40%",
                "Quality improvement: 30%",
                "",
                "🎓 CAMPUS RECRUITMENT",
                "Scenario: Recruiting 100 freshers",
                "Standardized evaluation",
                "Instant skill gap detection",
                "Personalized learning paths",
            ],
            [
                "🌐 FREELANCE MARKETPLACE",
                "Auto-match projects to freelancers",
                "Quick candidate evaluation",
                "90% reduction in matching time",
                "",
                "🔄 INTERNAL PROMOTIONS",
                "Objective skill assessment",
                "Clear improvement areas",
                "Training plan creation",
                "Readiness scoring",
                "",
                "📞 RECRUITMENT AGENCIES",
                "Scale candidate assessment",
                "Faster placement cycle",
                "Better candidate matching",
            ]
        )

        # SLIDE 16: ROI Analysis
        self.add_metric_slide(
            "Return on Investment (ROI) Analysis",
            [
                ("Year 1 Cost", "$15,900"),
                ("Annual Value", "$19,500+"),
                ("Net Return", "$3,600-10,100"),
                ("Break-even", "6 Months"),
            ]
        )

        # SLIDE 17: Implementation Timeline
        self.add_content_slide(
            "Implementation & Rollout",
            [
                "📅 WEEK 1-2: ONBOARDING",
                "✓ System setup & config",
                "✓ API key integration",
                "✓ Environment setup",
                "✓ Team training",
                "✓ Pilot testing",
                "",
                "📅 WEEK 3-4: ROLLOUT",
                "✓ Full deployment",
                "✓ User access setup",
                "✓ Sample data import",
                "✓ Documentation",
                "✓ Support training",
            ],
            [
                "📅 ONGOING: OPTIMIZATION",
                "✓ Performance monitoring",
                "✓ User feedback collection",
                "✓ Process refinement",
                "✓ Model improvements",
                "✓ Feature additions",
                "",
                "🚀 DEPLOYMENT OPTIONS",
                "Cloud: AWS, Azure, GCP",
                "On-premise: Docker container",
                "Hybrid: Mixed deployment",
                "",
                "📊 SUCCESS METRICS",
                "70% faster hiring",
                "40% cost reduction",
                "92% accuracy maintained",
            ]
        )

        # SLIDE 18: Roadmap & Future
        self.add_content_slide(
            "Roadmap & Future Enhancements",
            [
                "🎯 Q2 2026",
                "✓ Multi-language support",
                "✓ Video interview integration",
                "✓ Bias detection engine",
                "",
                "🚀 Q3 2026",
                "Social media profile analysis",
                "Mobile apps (iOS/Android)",
                "Advanced analytics dashboard",
                "",
                "🌍 Q4 2026",
                "Custom fine-tuned models",
                "Team collaboration features",
                "HRIS integration",
            ],
            [
                "📈 2027+ VISION",
                "Predictive success scoring",
                "AI interview coach",
                "Blockchain verification",
                "Global expansion",
                "",
                "💡 INNOVATION FOCUS",
                "Advanced bias mitigation",
                "Real-time video analysis",
                "Emotion recognition",
                "Personality assessment",
                "",
                "🔗 INTEGRATIONS",
                "ATS systems (Lever, Greenhouse)",
                "HRIS platforms",
                "Calendar & email",
                "Slack & Teams",
            ]
        )

        # SLIDE 19: Competitive Advantages
        self.add_content_slide(
            "Competitive Advantages",
            [
                "⚡ SPEED",
                "Resume analysis in 3-4 seconds",
                "No competitor matches this speed",
                "",
                "🎯 ACCURACY",
                "92%+ matching accuracy",
                "GPT-4o powered intelligence",
                "",
                "💡 INTELLIGENCE",
                "Semantic understanding",
                "Context-aware evaluation",
                "Personalized recommendations",
                "",
                "💰 VALUE",
                "60-70% time reduction",
                "40% cost reduction",
                "ROI in 6 months",
            ],
            [
                "🔐 SECURITY",
                "Enterprise-grade security",
                "JWT authentication",
                "Data encryption",
                "",
                "📊 ANALYTICS",
                "Real-time hiring metrics",
                "Data-driven insights",
                "Candidate comparisons",
                "",
                "🌐 SCALABILITY",
                "100+ concurrent users",
                "Handles millions of embeddings",
                "Cloud-ready architecture",
                "",
                "📱 EASY TO USE",
                "Intuitive Streamlit UI",
                "Minimal training required",
                "Fast implementation",
            ]
        )

        # SLIDE 20: Call to Action & Contact
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = PRIMARY_BLUE

        # Main message
        main_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
        main_frame = main_box.text_frame
        main_frame.text = "READY TO TRANSFORM YOUR HIRING?"
        main_frame.paragraphs[0].font.size = Pt(44)
        main_frame.paragraphs[0].font.bold = True
        main_frame.paragraphs[0].font.color.rgb = WHITE
        main_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # CTA
        cta_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(0.8))
        cta_frame = cta_box.text_frame
        cta_frame.text = "🚀 GET STARTED TODAY 🚀"
        cta_frame.paragraphs[0].font.size = Pt(32)
        cta_frame.paragraphs[0].font.bold = True
        cta_frame.paragraphs[0].font.color.rgb = ACCENT_GREEN
        cta_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Contact info
        contact_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(2))
        contact_frame = contact_box.text_frame
        contact_frame.word_wrap = True

        points = [
            "📧 Email: contact@recruitment-ai.com",
            "📱 Phone: +1-800-RECRUIT-AI",
            "🌐 Website: www.recruitment-ai.com",
        ]

        for i, point in enumerate(points):
            if i == 0:
                p = contact_frame.paragraphs[0]
            else:
                p = contact_frame.add_paragraph()
            p.text = point
            p.font.size = Pt(18)
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
            p.space_before = Pt(4)

    def save(self, filename):
        """Save the presentation"""
        self.prs.save(filename)
        print(f"✅ PowerPoint presentation created: {filename}")
        print(f"📊 Total slides: {len(self.prs.slides)}")

if __name__ == "__main__":
    print("🎨 Generating Recruitment AI System Presentation...")
    print("=" * 60)
    
    generator = PresentationGenerator()
    generator.generate()
    
    output_path = os.path.join(
        "d:/HR System  for the recruitment process/recruitment_ai",
        "Recruitment_AI_System_Presentation.pptx"
    )
    
    generator.save(output_path)
    
    print("=" * 60)
    print(f"✨ Presentation ready at:")
    print(f"   {output_path}")
    print()
    print("📋 Presentation includes:")
    print("   ✓ 20 professional slides")
    print("   ✓ Color-coded design (Blue/Green/Gray)")
    print("   ✓ All system information")
    print("   ✓ Architecture overview")
    print("   ✓ Performance metrics")
    print("   ✓ ROI analysis")
    print("   ✓ Implementation roadmap")
    print("   ✓ Call to action")
